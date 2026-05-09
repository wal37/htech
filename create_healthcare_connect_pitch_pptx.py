from __future__ import annotations

import sys
from pathlib import Path


BASE_DIR = Path(__file__).resolve().parent
for candidate in [
    Path("/Users/wale/Desktop/softtttt/aiapis/apiweb/.venv/lib/python3.13/site-packages"),
    Path("/Users/wale/Desktop/softtttt/foodup/foodweb/.venv/lib/python3.13/site-packages"),
]:
    if candidate.exists() and str(candidate) not in sys.path:
        sys.path.append(str(candidate))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = BASE_DIR / "HEALTHCARE_CONNECT_PITCH_DECK.pptx"


def rgb(code: str) -> RGBColor:
    code = code.replace("#", "")
    return RGBColor(int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))


PAPER = rgb("#EEF5F4")
WHITE = rgb("#FFFFFF")
INK = rgb("#0D2332")
SOFT = rgb("#607284")
LINE = rgb("#D4E1E8")
NAVY = rgb("#0F2E46")
COBALT = rgb("#1F5F93")
TEAL = rgb("#0C8A7B")
DEEP = rgb("#09151F")
DEEP_SOFT = rgb("#163149")
GOLD = rgb("#F1CF63")
MIST = rgb("#D8E6EE")


SLIDES = [
    {
        "type": "cover",
        "title": "Healthcare Connect",
        "subtitle": "A government-ready operating system for modern healthcare coordination, reporting, and AI-assisted oversight.",
        "eyebrow": "Health + Data + AI + Cloud",
        "stats": [
            ("$288.6B", "digital health"),
            ("$94.95B", "U.S. market"),
            ("$42.5B", "health cloud"),
        ],
    },
    {
        "type": "three_panel",
        "title": "Problem",
        "subtitle": "Most healthcare systems still run across fragmented records, slow routing, weak reporting, and disconnected oversight.",
        "items": [
            ("Fragmented care", "Patient records, referrals, service history, and follow-up actions are often split across facilities and siloed tools."),
            ("Weak visibility", "Ministries, hospital groups, and regional operators struggle to see demand, bottlenecks, and program performance in real time."),
            ("Data without action", "Health data may exist, but not inside a platform that supports secure workflows, routing, and governed decision-making."),
        ],
    },
    {
        "type": "signal",
        "title": "Solution",
        "subtitle": "Healthcare Connect turns fragmented health workflows into one secure cloud operating layer for care, reporting, and system oversight.",
        "bullets": [
            "Unified patient and case records across authorized health institutions",
            "Referral, scheduling, discharge, and follow-up coordination",
            "Facility, workforce, claims, and program reporting",
            "AI support for forecasting, triage, risk review, and anomaly detection",
        ],
    },
    {
        "type": "three_panel",
        "title": "Why now",
        "subtitle": "Three market forces make government-grade health infrastructure a high-conviction timing opportunity.",
        "items": [
            ("Digitization budgets are moving", "The global digital health market reached $288.55B in 2024 and is projected to reach $946.04B by 2030."),
            ("U.S. market is already large", "The U.S. digital health market is estimated at $94.95B in 2024 and projected to reach $276.62B by 2030."),
            ("AI + cloud need control", "Healthcare cloud reached $42.5B in 2024, while AI in healthcare is projected to reach $505.59B by 2033."),
        ],
    },
    {
        "type": "map",
        "title": "Product wedge",
        "subtitle": "Healthcare Connect sits at the operational center of care delivery, reporting, and health-system intelligence.",
        "left_title": "Inputs",
        "left_points": ["Patient data", "Facility activity", "Program reporting"],
        "center_title": "Healthcare Connect core",
        "center_points": ["Records", "Referrals", "Analytics", "AI support", "Oversight"],
        "right_title": "Outputs",
        "right_points": ["Care continuity", "Dashboards", "Forecasts", "Public reporting"],
    },
    {
        "type": "stack",
        "title": "Core modules",
        "subtitle": "The platform expands from a coordination wedge into broader system infrastructure over time.",
        "layers": [
            ("Patient record layer", "Longitudinal profiles, documents, diagnoses, medications, and encounter continuity."),
            ("Coordination layer", "Referrals, scheduling, case routing, care-team actions, and follow-up management."),
            ("Administration layer", "Claims, program eligibility, facility oversight, and workflow reporting."),
            ("Intelligence layer", "Dashboards, forecasting, anomaly detection, and governed AI support."),
        ],
    },
    {
        "type": "three_panel",
        "title": "Beachhead customers",
        "subtitle": "The first buyers are institutions with urgent coordination needs, budget authority, and recurring operational complexity.",
        "items": [
            ("National programs", "High-value deployments around immunization, maternal health, chronic care, disease programs, and public reporting."),
            ("Regional oversight", "Authorities need visibility across facilities, referrals, workforce activity, and health-system performance."),
            ("Provider networks", "Hospitals, labs, pharmacies, and regulated partners need one governed environment instead of disconnected tools."),
        ],
    },
    {
        "type": "timeline",
        "title": "How value is created",
        "subtitle": "The product produces operational value in a repeatable loop from intake to insight.",
        "steps": [
            ("Register", "Create or locate patient identity, authorized records, and service context."),
            ("Coordinate", "Route referrals, diagnostics, follow-up tasks, and cross-facility actions."),
            ("Report", "Aggregate facility, program, and regional activity into structured dashboards."),
            ("Improve", "Use analytics and AI support to plan capacity, flag risk, and strengthen performance."),
        ],
    },
    {
        "type": "ledger",
        "title": "Why this can win",
        "subtitle": "Trust, governance, and workflow depth create defensibility in regulated deployment environments.",
        "rows": [
            ("Access control", "Patient records and workflows", "defensible"),
            ("Audit trail", "Record activity and changes", "embedded"),
            ("Data protection", "Cloud and backup operations", "required"),
            ("Program oversight", "Regional reporting and review", "sticky"),
        ],
    },
    {
        "type": "three_panel",
        "title": "Market opportunity",
        "subtitle": "Healthcare Connect sits inside multiple large infrastructure markets with recurring modernization demand.",
        "items": [
            ("$288.55B digital health", "Global digital health was valued at $288.55B in 2024 and is projected to reach $946.04B by 2030."),
            ("$76.72B cloud infrastructure", "Global healthcare cloud infrastructure reached $76.72B in 2024, with 16.66% CAGR projected through 2030."),
            ("$505.59B AI in healthcare", "Grand View projects global AI in healthcare could reach $505.59B by 2033, with North America leading adoption."),
        ],
    },
    {
        "type": "signal",
        "title": "Go-to-market",
        "subtitle": "Commercial entry starts with high-need operational deployments and expands through network-level integration.",
        "bullets": [
            "Lead with high-friction coordination and reporting use cases",
            "Sell into ministries, regional authorities, and public hospital groups",
            "Expand from initial workflows into reporting, claims, and AI modules",
            "Build multi-year revenue through deployment, integration, and enterprise contracts",
        ],
    },
    {
        "type": "three_panel",
        "title": "Revenue model",
        "subtitle": "Revenue compounds through infrastructure contracts, implementation layers, and high-value add-on modules.",
        "items": [
            ("Platform licensing", "Government and enterprise health contracts tied to program, region, or network deployment."),
            ("Implementation revenue", "Integration, rollout, migration, training, and workflow configuration services."),
            ("Expansion modules", "Premium analytics, interoperability, claims workflows, and governed AI support."),
        ],
    },
    {
        "type": "timeline",
        "title": "Expansion roadmap",
        "subtitle": "The business expands from workflow infrastructure into broader system intelligence and network coverage.",
        "steps": [
            ("Phase 1", "Win with coordination, identity, referrals, and dashboards."),
            ("Phase 2", "Expand into claims, reporting, and population-level oversight."),
            ("Phase 3", "Layer in AI support, forecasting, and anomaly detection."),
            ("Phase 4", "Scale to national interoperability and program-specific extensions."),
        ],
    },
    {
        "type": "signal",
        "title": "Investment case",
        "subtitle": "Healthcare Connect sits where public-sector digitization, healthcare IT replacement, and governed AI adoption meet.",
        "bullets": [
            "Large multi-market opportunity across digital health, healthcare IT, AI, and cloud",
            "Sticky, workflow-embedded product with high switching costs once adopted",
            "Recurring enterprise revenue plus integration and expansion layers",
            "Strong fit for long-duration public-sector and regulated-health buyers",
        ],
    },
    {
        "type": "close",
        "title": "Healthcare Connect",
        "subtitle": "We are building the operating system for modern public-sector and regulated healthcare coordination.",
        "contact": "enquiry@healthcareconnect.ca",
    },
]


def add_textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_slide_number(slide, number: int):
    add_textbox(slide, Inches(10.25), Inches(5.9), Inches(0.45), Inches(0.18), f"{number:02d}", size=9, color=SOFT, font="Aptos", align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle, light=False):
    accent = GOLD if not light else TEAL
    title_color = WHITE if light else INK
    sub_color = MIST if light else SOFT
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(0.72), Inches(0.14), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, Inches(0.93), Inches(0.72), Inches(9.8), Inches(0.45), title, size=27, color=title_color, bold=True, font="Georgia")
    add_textbox(slide, Inches(0.93), Inches(1.23), Inches(9.7), Inches(0.56), subtitle, size=12, color=sub_color, font="Aptos")


def cover_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_card(slide, Inches(0.42), Inches(0.42), Inches(10.48), Inches(5.55), DEEP)
    add_card(slide, Inches(8.08), Inches(0.42), Inches(2.82), Inches(5.55), DEEP_SOFT)
    star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, Inches(8.72), Inches(0.95), Inches(1.48), Inches(1.48))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.color.rgb = GOLD
    add_textbox(slide, Inches(0.9), Inches(0.9), Inches(3.2), Inches(0.3), data["eyebrow"], size=13, color=rgb("#B8D3E6"), bold=True, font="Aptos")
    add_textbox(slide, Inches(0.88), Inches(1.52), Inches(6.2), Inches(0.95), data["title"], size=31, color=WHITE, bold=True, font="Georgia")
    add_textbox(slide, Inches(0.9), Inches(2.72), Inches(6.0), Inches(1.0), data["subtitle"], size=16, color=MIST, font="Aptos")
    add_textbox(slide, Inches(0.92), Inches(5.0), Inches(3.5), Inches(0.22), "enquiry@healthcareconnect.ca", size=12, color=MIST, bold=True, font="Aptos")
    for idx, (label, copy) in enumerate(data["stats"]):
        left = Inches(8.32)
        top = Inches(2.0 + idx * 1.16)
        add_card(slide, left, top, Inches(2.18), Inches(0.92), rgb("#1B3750"))
        add_textbox(slide, left + Inches(0.18), top + Inches(0.14), Inches(1.7), Inches(0.2), label, size=10, color=rgb("#A6C8DF"), bold=True, font="Aptos")
        add_textbox(slide, left + Inches(0.18), top + Inches(0.4), Inches(1.8), Inches(0.26), copy, size=12, color=WHITE, bold=True, font="Georgia")
    add_slide_number(slide, number)


def three_panel_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    for idx, (heading, copy) in enumerate(data["items"]):
        left = Inches(0.65 + idx * 3.55)
        add_card(slide, left, Inches(2.08), Inches(3.2), Inches(3.75), WHITE, LINE)
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, Inches(2.08), Inches(3.2), Inches(0.12))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD if idx != 1 else TEAL
        accent.line.color.rgb = GOLD if idx != 1 else TEAL
        add_textbox(slide, left + Inches(0.22), Inches(2.33), Inches(2.7), Inches(0.36), heading, size=18, color=NAVY, bold=True, font="Georgia")
        add_textbox(slide, left + Inches(0.22), Inches(2.92), Inches(2.75), Inches(1.52), copy, size=12, color=SOFT, font="Aptos")
    add_slide_number(slide, number)


def signal_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_card(slide, Inches(0.52), Inches(0.52), Inches(10.3), Inches(5.9), PAPER)
    add_title(slide, data["title"], data["subtitle"])
    for idx, bullet in enumerate(data["bullets"]):
        y = Inches(2.16 + idx * 0.92)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95), y + Inches(0.08), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.color.rgb = GOLD
        add_textbox(slide, Inches(1.28), y, Inches(8.7), Inches(0.3), bullet, size=16, color=INK, font="Aptos")
    add_slide_number(slide, number)


def map_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    columns = [
        (Inches(0.85), Inches(2.08), Inches(2.55), Inches(3.55), WHITE, data["left_title"], data["left_points"]),
        (Inches(4.1), Inches(1.82), Inches(3.25), Inches(4.05), NAVY, data["center_title"], data["center_points"]),
        (Inches(8.05), Inches(2.08), Inches(2.55), Inches(3.55), WHITE, data["right_title"], data["right_points"]),
    ]
    for left, top, width, height, fill, title, points in columns:
        add_card(slide, left, top, width, height, fill, LINE if fill == WHITE else NAVY)
        title_color = WHITE if fill == NAVY else INK
        body_color = MIST if fill == NAVY else SOFT
        add_textbox(slide, left + Inches(0.22), top + Inches(0.2), width - Inches(0.44), Inches(0.35), title, size=18, color=title_color, bold=True, font="Georgia")
        for idx, point in enumerate(points):
            add_textbox(slide, left + Inches(0.28), top + Inches(0.84 + idx * 0.58), width - Inches(0.56), Inches(0.24), f"• {point}", size=12, color=body_color, font="Aptos")
    add_slide_number(slide, number)


def stack_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, data["title"], data["subtitle"])
    fills = [rgb("#EAF2F6"), rgb("#DDEBF3"), rgb("#EAF2F6"), rgb("#DDEBF3")]
    for idx, (heading, copy) in enumerate(data["layers"]):
        top = Inches(2.0 + idx * 1.0)
        add_card(slide, Inches(0.78), top, Inches(9.95), Inches(0.78), fills[idx], rgb("#C9DDE7"))
        add_textbox(slide, Inches(1.02), top + Inches(0.16), Inches(2.4), Inches(0.22), heading, size=15, color=NAVY, bold=True, font="Georgia")
        add_textbox(slide, Inches(3.2), top + Inches(0.14), Inches(7.1), Inches(0.26), copy, size=11, color=INK, font="Aptos")
    add_slide_number(slide, number)


def timeline_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP
    add_title(slide, data["title"], data["subtitle"], light=True)
    for idx, (heading, copy) in enumerate(data["steps"]):
        left = Inches(0.82 + idx * 2.48)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.58), Inches(2.1), Inches(2.48))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb("#173044")
        shape.line.color.rgb = rgb("#22445F")
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), Inches(2.82), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GOLD
        badge.line.color.rgb = GOLD
        add_textbox(slide, left + Inches(0.62), Inches(2.82), Inches(0.4), Inches(0.18), f"{idx + 1:02d}", size=11, color=rgb("#A7C9DF"), bold=True, font="Aptos")
        add_textbox(slide, left + Inches(0.18), Inches(3.25), Inches(1.72), Inches(0.28), heading, size=17, color=WHITE, bold=True, font="Georgia")
        add_textbox(slide, left + Inches(0.18), Inches(3.8), Inches(1.72), Inches(0.78), copy, size=11, color=MIST, font="Aptos")
    add_slide_number(slide, number)


def ledger_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    add_card(slide, Inches(0.72), Inches(2.0), Inches(10.0), Inches(3.82), WHITE, LINE)
    add_textbox(slide, Inches(1.0), Inches(2.24), Inches(2.3), Inches(0.18), "Control", size=11, color=SOFT, bold=True, font="Aptos")
    add_textbox(slide, Inches(4.4), Inches(2.24), Inches(2.8), Inches(0.18), "Record or workflow", size=11, color=SOFT, bold=True, font="Aptos")
    add_textbox(slide, Inches(8.92), Inches(2.24), Inches(1.1), Inches(0.18), "Status", size=11, color=SOFT, bold=True, font="Aptos")
    for idx, row in enumerate(data["rows"]):
        y = Inches(2.74 + idx * 0.72)
        rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.94), y - Inches(0.06), Inches(9.55), Inches(0.01))
        rule.fill.solid()
        rule.fill.fore_color.rgb = rgb("#D9E7EE")
        rule.line.color.rgb = rgb("#D9E7EE")
        add_textbox(slide, Inches(1.0), y, Inches(2.85), Inches(0.22), row[0], size=12, color=INK, font="Aptos")
        add_textbox(slide, Inches(4.4), y, Inches(3.1), Inches(0.22), row[1], size=12, color=SOFT, font="Aptos")
        add_textbox(slide, Inches(8.98), y, Inches(1.0), Inches(0.2), row[2], size=10, color=TEAL, bold=True, font="Aptos")
    add_slide_number(slide, number)


def close_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP
    add_textbox(slide, Inches(0.9), Inches(1.08), Inches(2.8), Inches(0.26), "Healthcare Infrastructure", size=14, color=rgb("#B8D3E6"), bold=True, font="Aptos")
    add_textbox(slide, Inches(0.88), Inches(1.88), Inches(5.8), Inches(0.8), data["title"], size=30, color=WHITE, bold=True, font="Georgia")
    add_textbox(slide, Inches(0.9), Inches(2.88), Inches(6.4), Inches(0.95), data["subtitle"], size=16, color=MIST, font="Aptos")
    add_textbox(slide, Inches(0.92), Inches(4.95), Inches(4.0), Inches(0.22), data["contact"], size=14, color=WHITE, bold=True, font="Aptos")
    add_card(slide, Inches(7.92), Inches(1.0), Inches(2.42), Inches(4.92), DEEP_SOFT)
    star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, Inches(8.43), Inches(1.6), Inches(1.4), Inches(1.4))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.color.rgb = GOLD
    add_textbox(slide, Inches(8.12), Inches(3.64), Inches(2.0), Inches(0.24), "Government-ready", size=11, color=MIST, bold=True, font="Aptos", align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.12), Inches(4.0), Inches(2.0), Inches(0.48), "care, data,\nand AI", size=15, color=WHITE, bold=True, font="Georgia", align=PP_ALIGN.CENTER)
    add_slide_number(slide, number)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(11.33)
    prs.slide_height = Inches(6.38)

    for idx, slide_data in enumerate(SLIDES, start=1):
        kind = slide_data["type"]
        if kind == "cover":
            cover_slide(prs, slide_data, idx)
        elif kind == "three_panel":
            three_panel_slide(prs, slide_data, idx)
        elif kind == "signal":
            signal_slide(prs, slide_data, idx)
        elif kind == "map":
            map_slide(prs, slide_data, idx)
        elif kind == "stack":
            stack_slide(prs, slide_data, idx)
        elif kind == "timeline":
            timeline_slide(prs, slide_data, idx)
        elif kind == "ledger":
            ledger_slide(prs, slide_data, idx)
        elif kind == "close":
            close_slide(prs, slide_data, idx)

    prs.core_properties.title = "Healthcare Connect Pitch Deck"
    prs.core_properties.subject = "Government-ready digital health infrastructure"
    prs.core_properties.author = "OpenAI Codex"
    prs.save(OUT)


if __name__ == "__main__":
    build_deck()
    print(OUT)
